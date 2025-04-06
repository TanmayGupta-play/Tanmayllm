import re
import os
import random
from flask import Flask, request, jsonify, send_file, render_template
from flask_cors import CORS  # Import for handling CORS
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import MSO_AUTO_SIZE, PP_ALIGN
from pptx.dml.color import RGBColor
import gptText
import addphoto

app = Flask(__name__)
CORS(app)  # Enable CORS for all routes

# Define upload and output directories
UPLOAD_FOLDER = 'uploads'
OUTPUT_FOLDER = 'static/presentations'
os.makedirs(UPLOAD_FOLDER, exist_ok=True)
os.makedirs(OUTPUT_FOLDER, exist_ok=True)

# Define template paths
TEMPLATES = {
    1: "template/minimalistic.pptx",
    2: "template/colourful.pptx",
    3: "template/professional.pptx",
    4: "template/dark.pptx",
}

# Store the latest generated presentation for the /result endpoint
current_presentation = None

def gettext(topic_list, code: bool):
    slides_data = gptText.structured(topic_list=topic_list, include_code=code)
    if slides_data and isinstance(slides_data, list) and len(slides_data) > 0 and 'Slides' in slides_data[0]:
        return slides_data[0]['Slides']
    else:
        return ""

def getphoto(slide_data):
    os.makedirs('images', exist_ok=True)
    pattern = r"Image Suggestion:\s*(.+)"
    image_suggestions = re.findall(pattern, slide_data)
    image_paths = []
    for suggestion in image_suggestions:
        try:
            image_result = addphoto.get_images(suggestion, 1)
            if image_result:
                image_paths.append(image_result[0])
            else:
                print(f"No image found for suggestion: {suggestion}")
        except Exception as e:
            print(f"Error getting image for suggestion '{suggestion}': {e}")
    return image_paths

def create_presentation(topic, template_choice=1, include_code=False):
    """Create a presentation based on topic and template choice"""
    # Convert topic to list if it's a string
    topic_list = [topic] if isinstance(topic, str) else topic
    
    # Generate a unique filename for the presentation
    presentation_id = f"presentation_{random.randint(1000, 9999)}"
    output_filename = f"{OUTPUT_FOLDER}/{presentation_id}.pptx"
    
    # Get slide data and image paths
    slide_data = gettext(topic_list, include_code)
    image_paths = getphoto(slide_data)

    if image_paths is None:
        image_paths = []

    # Get template path
    template_path = TEMPLATES.get(template_choice, None)
    
    # Load the selected template or create a blank presentation
    if template_path and os.path.exists(template_path):
        prs = Presentation(template_path)
        # Clear all existing slides
        xml_slides = prs.slides._sldIdLst  # Access the slide ID list
        slides = list(xml_slides)
        for slide in slides:
            xml_slides.remove(slide)  # Remove each slide
    else:
        prs = Presentation()

    # Add a blank slide layout to start fresh
    blank_slide_layout = prs.slide_layouts[6]

    # Extract the title from the slide data
    title_match = re.search(r"Title:\s*(.*)", slide_data)
    if title_match:
        overall_title = title_match.group(1).strip()
        title_slide_layout = prs.slide_layouts[0]
        slide = prs.slides.add_slide(title_slide_layout)

        # Check if the title placeholder exists
        if slide.shapes.title:
            title = slide.shapes.title
            title.text = overall_title
        else:
            # Add a new text box for the title if the placeholder doesn't exist
            title_box = slide.shapes.add_textbox(Inches(0.25), Inches(0.5), Inches(9), Inches(1.5))
            title_frame = title_box.text_frame
            title_frame.text = overall_title
            title_p = title_frame.paragraphs[0]
            title_p.font.bold = True
            title_p.font.size = Pt(32)
            title_p.alignment = PP_ALIGN.CENTER

        # Check if a subtitle placeholder exists
        if len(slide.placeholders) > 1:
            subtitle = slide.placeholders[1]
            subtitle.text = "Generated Presentation"

    # Split the slide data into individual slides
    slides = re.split(r"---", slide_data)
    image_index = 0

    for slide_content in slides:
        slide_content = slide_content.strip()
        if not slide_content:
            continue

        # Extract the slide title
        title_match = re.search(r"Slide\s*\d+:\s*(.*)", slide_content)
        slide_title = title_match.group(1).strip() if title_match else "Slide"

        # Add a new slide using the blank layout
        slide = prs.slides.add_slide(blank_slide_layout)

        # Add the title to the slide
        title_box = slide.shapes.add_textbox(Inches(1), Inches(0.3), Inches(8), Inches(1))
        title_frame = title_box.text_frame
        title_frame.text = slide_title
        title_p = title_frame.paragraphs[0]
        title_p.font.bold = True
        title_p.font.size = Pt(32)
        title_p.alignment = PP_ALIGN.CENTER

        # Add the content to the slide
        content_box = slide.shapes.add_textbox(Inches(0.5), Inches(1), Inches(6), Inches(5))
        tf = content_box.text_frame
        tf.word_wrap = True
        tf.auto_size = MSO_AUTO_SIZE.SHAPE_TO_FIT_TEXT

        lines = slide_content.split("\n")
        for line in lines:
            line = line.strip()

            if line.startswith("Slide") and ":" in line:
                continue
            elif line.startswith("Image Suggestion:"):
                continue
            elif line.startswith("-"):
                p = tf.add_paragraph()
                p.text = line[1:].strip()
                p.level = 0
                p.font.size = Pt(20)
            elif line.startswith("```"):
                code_match = re.search(r"```[a-zA-Z]*\n(.*?)```", slide_content, re.DOTALL)
                if code_match:
                    code = code_match.group(1).strip()
                    code_slide = prs.slides.add_slide(blank_slide_layout)
                    code_title_box = code_slide.shapes.add_textbox(Inches(1), Inches(0.3), Inches(8), Inches(1))
                    code_title_frame = code_title_box.text_frame
                    code_title_frame.text = "Code Example"
                    code_title_p = code_title_frame.paragraphs[0]
                    code_title_p.font.bold = True
                    code_title_p.font.size = Pt(32)
                    code_title_p.alignment = PP_ALIGN.CENTER

                    code_box = code_slide.shapes.add_textbox(Inches(1), Inches(1.5), Inches(8), Inches(5))
                    code_tf = code_box.text_frame
                    code_tf.word_wrap = True
                    code_tf.auto_size = MSO_AUTO_SIZE.SHAPE_TO_FIT_TEXT
                    p = code_tf.add_paragraph()
                    p.text = code
                    p.font.name = 'Courier New'
                    p.font.size = Pt(10)
                    p.font.color.rgb = RGBColor(255, 255, 255)
                    p.font.fill.solid()
                    p.font.fill.fore_color.rgb = RGBColor(40, 40, 40)
            elif line:
                p = tf.add_paragraph()
                p.text = line
                p.font.size = Pt(20)

        # Add an image to the slide if available
        if image_paths and image_index < len(image_paths):
            try:
                img_path = image_paths[image_index]

                # Define the dimensions for the image
                left = Inches(6.5)  # Place the image on the right side
                top = Inches(1)  # Align the image with the top of the content
                width = Inches(3)  # Elongated width
                height = Inches(4)  # Elongated height

                # Add the image to the slide
                slide.shapes.add_picture(img_path, left, top, width, height)
            except FileNotFoundError:
                print(f"Image file not found: {img_path}")
            except Exception as e:
                print(f"Error adding image: {e}")
            image_index += 1

    # Save the presentation
    prs.save(output_filename)
    print(f"Presentation saved to {output_filename}")
    
    # Clean up the images directory
    if os.path.exists('images'):
        try:
            for filename in os.listdir('images'):
                file_path = os.path.join('images', filename)
                os.remove(file_path)
            os.rmdir('images')
        except OSError as e:
            print(f"Failed to cleanup 'images' directory: {e}")
    
    return {
        "id": presentation_id,
        "filename": f"{presentation_id}.pptx",
        "path": output_filename
    }

# API Routes for React Frontend

@app.route('/api/generate', methods=['POST'])
def api_generate():
    """API endpoint for generating presentations"""
    global current_presentation
    
    try:
        data = request.get_json()
        
        if not data or 'topic' not in data:
            return jsonify({"error": "Topic is required"}), 400
        
        topic = data.get('topic')
        template = int(data.get('template', 1))
        include_code = data.get('includeCode', False)
        
        # Generate presentation
        presentation_info = create_presentation(topic, template, include_code)
        current_presentation = presentation_info
        
        return jsonify({
            "message": "Presentation created successfully!",
            "presentation_id": presentation_info["id"],
            "filename": presentation_info["filename"]
        }), 200
    
    except Exception as e:
        print(f"Error generating presentation: {str(e)}")
        return jsonify({"error": str(e)}), 500

@app.route('/api/result', methods=['GET'])
def api_result():
    """API endpoint to get information about the current presentation"""
    global current_presentation
    
    if current_presentation:
        return jsonify({
            "presentation_id": current_presentation["id"],
            "filename": current_presentation["filename"]
        }), 200
    else:
        return jsonify({"error": "No presentation has been generated yet"}), 404

@app.route('/api/download/<presentation_id>', methods=['GET'])
def api_download(presentation_id):
    """API endpoint to download a generated presentation"""
    filename = f"{presentation_id}.pptx"
    file_path = os.path.join(OUTPUT_FOLDER, filename)
    
    if os.path.exists(file_path):
        return send_file(file_path, as_attachment=True, download_name=filename)
    else:
        return jsonify({"error": "Presentation not found"}), 404

@app.route('/')
def index():
    """Render the main page - for direct browser access"""
    templates = [{"id": key, "name": os.path.basename(path).replace(".pptx", "").title()} 
                for key, path in TEMPLATES.items()]
    return render_template('index.html', templates=templates)

if __name__ == '__main__':
    app.run(debug=True)