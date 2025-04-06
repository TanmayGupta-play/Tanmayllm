import os
import re
import random
from flask import Flask, request, send_file, jsonify
from flask_cors import CORS
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import MSO_AUTO_SIZE, PP_ALIGN
from pptx.dml.color import RGBColor
import gptText
import addphoto

app = Flask(__name__)
CORS(app, resources={r"/api/*": {"origins": "*"}})  # Allow CORS for /api routes

# ✅ Define template paths
TEMPLATES = {
    1: "template/minimalistic.pptx",
    2: "template/colourful.pptx",
    3: "template/professional.pptx",
    4: "template/dark.pptx",
}


def gettext(topicList, code: bool):
    slides_data = gptText.structured(topic_list=topicList, include_code=code)
    if slides_data and isinstance(slides_data, list) and len(slides_data) > 0 and 'Slides' in slides_data[0]:
        return slides_data[0]['Slides']
    else:
        return ""

def getphoto(SlideData):
    pattern = r"Image Suggestion:\s*(.+)"
    image_suggestions = re.findall(pattern, SlideData)
    image_paths = []
    for suggestion in image_suggestions:
        try:
            image_result = addphoto.get_images(suggestion, 1)
            if image_result:
                image_paths.append(image_result[0])
            else:
                print(f"No image found for suggestion: {suggestion}")
        except Exception as e:
            print(f"Error getting image for suggestion '{suggestion}': {e}")
    return image_paths


def create_presentation(slide_data, image_paths=None, output_filename="presentation.pptx", template_path=None):
    """Create a PowerPoint presentation from slide data."""
    # Get slide data and image paths

    if image_paths is None:
        image_paths = []

    # Load the selected template or create a blank presentation
    if template_path and os.path.exists(template_path):
        prs = Presentation(template_path)
        # Clear all existing slides
        xml_slides = prs.slides._sldIdLst
        slides = list(xml_slides)
        for slide in slides:
            xml_slides.remove(slide)
    else:
        prs = Presentation()

    # Add a blank slide layout to start fresh
    blank_slide_layout = prs.slide_layouts[6]

    # Extract the title from the slide data
    title_match = re.search(r"Title:\s*(.*)", slide_data)
    if title_match:
        overall_title = title_match.group(1).strip()
        title_slide_layout = prs.slide_layouts[0]
        slide = prs.slides.add_slide(title_slide_layout)

        # Add title if available
        if slide.shapes.title:
            title = slide.shapes.title
            title.text = overall_title
        else:
            title_box = slide.shapes.add_textbox(Inches(1), Inches(0.5), Inches(5), Inches(1))
            title_frame = title_box.text_frame
            title_frame.text = overall_title
            title_p = title_frame.paragraphs[0]
            title_p.font.bold = True
            title_p.font.size = Pt(32)
            title_p.alignment = PP_ALIGN.CENTER

    slides = re.split(r"---", slide_data)
    image_index = 0

    for slide_content in slides:
        slide_content = slide_content.strip()
        if not slide_content:
            continue

        # Extract slide title
        title_match = re.search(r"Slide\s*\d+:\s*(.*)", slide_content)
        slide_title = title_match.group(1).strip() if title_match else "Slide"

        # Add a new slide using blank layout
        slide = prs.slides.add_slide(blank_slide_layout)

        # Add title
        title_box = slide.shapes.add_textbox(Inches(1), Inches(0.3), Inches(8), Inches(1))
        title_frame = title_box.text_frame
        title_frame.text = slide_title
        title_p = title_frame.paragraphs[0]
        title_p.font.bold = True
        title_p.font.size = Pt(32)
        title_p.alignment = PP_ALIGN.CENTER

        # Add slide content
        content_box = slide.shapes.add_textbox(Inches(0.5), Inches(1), Inches(6), Inches(5))
        tf = content_box.text_frame
        tf.word_wrap = True
        tf.auto_size = MSO_AUTO_SIZE.SHAPE_TO_FIT_TEXT

        lines = slide_content.split("\n")
        for line in lines:
            line = line.strip()
            if line.startswith(("Slide", "Image Suggestion:")):
                continue
            elif line:
                p = tf.add_paragraph()
                p.text = line
                p.font.size = Pt(20)

        # Add image to slide if available
        if image_paths and image_index < len(image_paths):
            try:
                img_path = image_paths[image_index]
                if os.path.exists(img_path):
                    left, top, width, height = Inches(6.5), Inches(1), Inches(3), Inches(5)
                    slide.shapes.add_picture(img_path, left, top, width, height)
                else:
                    print(f"Image not found: {img_path}")
            except Exception as e:
                print(f"Error adding image: {e}")
            image_index += 1

    # Save the presentation
    prs.save(output_filename)
    print(f"Presentation saved to {output_filename}")

    # Cleanup images if directory exists
    cleanup_images()

    return output_filename


def cleanup_images():
    """Remove generated images after presentation creation."""
    images_dir = "images"
    if os.path.exists(images_dir):
        try:
            for filename in os.listdir(images_dir):
                file_path = os.path.join(images_dir, filename)
                if os.path.isfile(file_path):
                    os.remove(file_path)
            os.rmdir(images_dir)
            print("✅ Image directory cleaned successfully.")
        except OSError as e:
            print(f"⚠️ Failed to clean up 'images' directory: {e}")


@app.route("/api/generate", methods=["POST"])
def generate_presentation():
    """Generate and return a presentation."""
    try:
        data = request.json
        topic = data.get("topic")
        template_choice = int(data.get("template", 1))
        template_path = TEMPLATES.get(template_choice, None)

        if not topic:
            return jsonify({"error": "Topic is required"}), 400

        # Generate slide data and images
        slide_data = gettext([topic], True)
        image_paths = getphoto(slide_data)

        # Create the presentation
        output_filename = "presentation.pptx"
        create_presentation(slide_data, image_paths, output_filename, template_path)

        return jsonify({"message": "Presentation created successfully!", "download_url": "/api/download"})

    except Exception as e:
        print(f"⚠️ Error in generating presentation: {e}")
        return jsonify({"error": "Failed to generate presentation", "details": str(e)}), 500


@app.route("/api/download", methods=["GET"])
def download_presentation():
    """Download the generated PPTX."""
    try:
        if os.path.exists("presentation.pptx"):
            return send_file("presentation.pptx", as_attachment=True)
        else:
            return jsonify({"error": "Presentation file not found"}), 404
    except Exception as e:
        return jsonify({"error": "Error downloading presentation", "details": str(e)}), 500


@app.route("/api/health", methods=["GET"])
def health_check():
    """API health check endpoint."""
    return jsonify({"status": "API is running!"}), 200


if __name__ == "__main__":
    app.run(debug=True, host="0.0.0.0", port=5000)
